from flask import Flask, render_template_string, request, redirect, url_for, session
import os
import tempfile
from ia_engine import generer_article
import PyPDF2
import docx
import pptx
import pandas as pd
import requests
from bs4 import BeautifulSoup
import html

from flask import session, redirect, url_for, flash
app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'dev-secret-key')  # À personnaliser en prod


def parse_id_list(raw_value: str) -> list[int]:
    ids = []
    for token in (raw_value or '').replace(';', ',').split(','):
        cleaned = token.strip()
        if cleaned.isdigit():
            ids.append(int(cleaned))
    return ids


def get_wp_auth():
    wp_url = os.environ.get('WP_URL', '').strip().rstrip('/')
    wp_user = os.environ.get('WP_USERNAME', '').strip()
    wp_password = os.environ.get('WP_APP_PASSWORD', '').strip()
    return wp_url, wp_user, wp_password


def _wp_auth_variants(wp_user: str, wp_password: str) -> list[tuple[str, str]]:
    """Retourne des variantes d'auth pour gérer les app passwords copiés avec espaces."""
    variants: list[tuple[str, str]] = []
    first = (wp_user, wp_password)
    variants.append(first)

    compact_password = ''.join(wp_password.split())
    compact = (wp_user, compact_password)
    if compact_password and compact != first:
        variants.append(compact)
    return variants


def wp_request(method: str, url: str, wp_user: str, wp_password: str, **kwargs) -> requests.Response:
    """Fait une requête WP avec retry auto sur variante sans espaces du mot de passe."""
    last_resp = None
    for auth_user, auth_password in _wp_auth_variants(wp_user, wp_password):
        resp = requests.request(method, url, auth=(auth_user, auth_password), **kwargs)
        last_resp = resp
        if resp.status_code != 401:
            return resp
    return last_resp


def get_wp_post_type() -> str:
    raw = (os.environ.get('WP_POST_TYPE', 'posts') or 'posts').strip().lower()
    # Autorise seulement les caractères valides d'un slug WP REST.
    safe = ''.join(ch for ch in raw if ch.isalnum() or ch in ('-', '_'))
    return safe or 'posts'


def _wp_body_snippet(resp: requests.Response, limit: int = 500) -> str:
    raw = (resp.text or '').strip()
    if not raw:
        return 'réponse vide'
    if len(raw) > limit:
        raw = raw[:limit] + '...'
    return html.escape(raw)


def build_wp_error_message(resp: requests.Response) -> str:
    """Construit un message d'erreur WordPress lisible et actionnable."""
    default = f"Erreur WordPress ({resp.status_code}) sur {html.escape(resp.url)}: {_wp_body_snippet(resp)}"

    if resp.status_code >= 500:
        server = html.escape(resp.headers.get('server', 'inconnu'))
        return (
            f"Erreur WordPress ({resp.status_code}) sur {html.escape(resp.url)}. "
            "Le serveur WordPress a planté avant de répondre correctement à l'API REST. "
            "Causes fréquentes: .htaccess invalide, plugin sécurité/cache, erreur PHP. "
            f"Serveur détecté: {server}. "
            f"Extrait réponse: {_wp_body_snippet(resp)}"
        )

    try:
        data = resp.json()
    except Exception:
        return default

    code = data.get('code', '')
    message = data.get('message', '')
    status = (data.get('data') or {}).get('status', resp.status_code)

    if code == 'rest_cannot_create' and status in (401, 403):
        post_type = get_wp_post_type()
        return (
            "WordPress refuse la création d'article (rest_cannot_create). "
            f"L'utilisateur authentifié n'a pas le droit de créer des contenus de type '{post_type}'. "
            "Vérifiez que WP_USERNAME a le rôle Auteur/Éditeur/Administrateur, "
            "que le mot de passe d'application est bien créé pour CE même utilisateur, "
            "et que les extensions de sécurité ne bloquent pas l'API REST. "
            f"Réponse WP: {html.escape(str(message))}"
        )

    if code == 'rest_not_logged_in' and status == 401:
        return (
            "WordPress ne reçoit pas une authentification valide (rest_not_logged_in). "
            "Vérifiez WP_USERNAME/WP_APP_PASSWORD (sans guillemets), "
            "puis côté hébergeur autorisez le header Authorization pour l'API REST. "
            "Sur Apache/LiteSpeed, ajoutez dans .htaccess: "
            "RewriteRule .* - [E=HTTP_AUTHORIZATION:%{HTTP:Authorization}] "
            f"Réponse WP: {html.escape(str(message))}"
        )

    return f"Erreur WordPress ({status}) [{html.escape(str(code))}]: {html.escape(str(message))}"

LOGIN_FORM = '''
<!DOCTYPE html>
<html lang="fr">
<head><meta charset="UTF-8"><title>Connexion</title></head>
<body>
<h2>Connexion requise</h2>
<form method="post">
    <input type="password" name="password" placeholder="Mot de passe" required autofocus>
    <button type="submit">Se connecter</button>
    {% if error %}<div style="color:red;">{{ error }}</div>{% endif %}
</form>
</body></html>
'''

HTML_FORM = '''
<!DOCTYPE html>
<html lang="fr">
<head>
    <meta charset="UTF-8">
    <title>Générateur d'articles IA</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; background: #f8f9fa; }
        textarea, input[type="file"] { width: 100%; margin-bottom: 10px; }
        textarea { height: 80px; border-radius: 6px; border: 1px solid #bbb; padding: 8px; }
        .result { margin-top: 30px; padding: 20px; background: #f4f4f4; border-radius: 8px; }
        .btn-main, .btn-redesign {
            background: linear-gradient(90deg, #007bff 60%, #0056b3 100%);
            color: #fff; border: none; border-radius: 6px; padding: 12px 28px;
            font-size: 1.1em; font-weight: bold; cursor: pointer; transition: background 0.2s, box-shadow 0.2s;
            box-shadow: 0 2px 8px #cce3ff;
        }
        .btn-main:hover, .btn-redesign:hover {
            background: linear-gradient(90deg, #0056b3 60%, #007bff 100%);
            box-shadow: 0 4px 16px #b3d1ff;
        }
        .file-list { margin: 8px 0 16px 0; padding: 8px; background: #e9ecef; border-radius: 6px; font-size: 0.98em; }
        .file-list span { display: inline-block; margin-right: 10px; }
        .feedback { margin: 10px 0; color: #007bff; font-size: 0.98em; }
    </style>
</head>
<body>
        <h1>Générateur d'articles IA</h1>
        <form method="post" id="mainForm" enctype="multipart/form-data">
            <label>Vos idées :</label><br>
            <textarea name="ideas" id="ideas" required>{{ ideas or '' }}</textarea><br>
            <label>Prompt de recherche :</label><br>
            <textarea name="prompt" id="prompt" required>{{ prompt or '' }}</textarea><br>
            <label>Joindre des fichiers (PDF, DOCX, TXT, etc.) :</label><br>
            <div style="position:relative;">
                <input type="file" id="documents" name="documents" multiple style="opacity:0;position:absolute;left:0;top:0;width:100%;height:40px;z-index:2;cursor:pointer;">
                <button type="button" id="customUploadBtn" class="btn-main" style="width:100%;margin-bottom:10px;z-index:1;position:relative;">Sélectionner des fichiers</button>
            </div>
            <div class="file-list" id="fileList"></div>
            <label>Ajouter des liens web (un par ligne) :</label><br>
            <textarea name="links" id="links" style="height:60px;" placeholder="https://exemple.com/article1\nhttps://exemple.com/article2"></textarea><br>
            <button type="submit" class="btn-main">Générer l'article</button>
        </form>
        <script>
        // Sauvegarde et restauration des champs dans localStorage
        const ideasField = document.getElementById('ideas');
        const promptField = document.getElementById('prompt');
        const linksField = document.getElementById('links');
        if(localStorage.getItem('ideas')) ideasField.value = localStorage.getItem('ideas');
        if(localStorage.getItem('prompt')) promptField.value = localStorage.getItem('prompt');
        if(localStorage.getItem('links')) linksField.value = localStorage.getItem('links');
        ideasField.addEventListener('input', () => localStorage.setItem('ideas', ideasField.value));
        promptField.addEventListener('input', () => localStorage.setItem('prompt', promptField.value));
        linksField.addEventListener('input', () => localStorage.setItem('links', linksField.value));

        // Upload moderne avec suppression possible
        const fileInput = document.getElementById('documents');
        const fileList = document.getElementById('fileList');
        const customUploadBtn = document.getElementById('customUploadBtn');
        let filesArray = [];

        customUploadBtn.addEventListener('click', () => fileInput.click());

        fileInput.addEventListener('change', (e) => {
            filesArray = Array.from(fileInput.files);
            renderFileList();
        });

        function renderFileList() {
            fileList.innerHTML = '';
            if (filesArray.length === 0) {
                fileList.innerHTML = '<span style="color:#888;">Aucun fichier sélectionné</span>';
                fileInput.value = '';
                return;
            }
            filesArray.forEach((file, idx) => {
                const fileItem = document.createElement('div');
                fileItem.style.display = 'flex';
                fileItem.style.alignItems = 'center';
                fileItem.style.marginBottom = '6px';
                fileItem.innerHTML = `<span style="background:#d1e7dd;padding:4px 10px;border-radius:5px;display:inline-block;margin-right:10px;">${file.name}</span>`;
                const removeBtn = document.createElement('button');
                removeBtn.textContent = '✖';
                removeBtn.type = 'button';
                removeBtn.style.background = '#ffdddd';
                removeBtn.style.color = '#c00';
                removeBtn.style.border = 'none';
                removeBtn.style.borderRadius = '50%';
                removeBtn.style.width = '28px';
                removeBtn.style.height = '28px';
                removeBtn.style.fontWeight = 'bold';
                removeBtn.style.cursor = 'pointer';
                removeBtn.style.marginLeft = '5px';
                removeBtn.addEventListener('click', () => {
                    filesArray.splice(idx, 1);
                    renderFileList();
                });
                fileItem.appendChild(removeBtn);
                fileList.appendChild(fileItem);
            });
        }

        // Avant soumission, reconstruire l'objet FileList à partir de filesArray
        document.getElementById('mainForm').addEventListener('submit', function() {
            try {
                if (fileInput.files.length !== filesArray.length && typeof DataTransfer !== 'undefined') {
                    // Reconstruit les fichiers sans bloquer la soumission si le navigateur ne supporte pas DataTransfer.
                    const dt = new DataTransfer();
                    filesArray.forEach(f => dt.items.add(f));
                    fileInput.files = dt.files;
                }
            } catch (err) {
                console.warn('Soumission sans synchronisation filesArray:', err);
            }
        });
        // Affichage initial
        renderFileList();
        </script>
        {% if feedback_links %}
        <div class="feedback">
            <strong>Résultat extraction des liens :</strong><br>
            {% for line in feedback_links.split('<br>') if line.strip() %}
                {% if 'Succès extraction' in line %}
                    <span style="color:green;font-weight:bold;">✔</span> {{ line|safe }}<br>
                {% elif 'Échec extraction' in line %}
                    <span style="color:red;font-weight:bold;">✖</span> {{ line|safe }}<br>
                {% else %}
                    {{ line|safe }}<br>
                {% endif %}
            {% endfor %}
        </div>
        {% endif %}
        {% if generation_error %}
        <div style="margin:10px 0;padding:10px;border-radius:6px;background:#f8d7da;color:#721c24;">
            {{ generation_error }}
        </div>
        {% endif %}
        {% if html_output %}
        <div class="result">
            <h2>Aperçu visuel WordPress :</h2>
            <div id="wp-preview" style="background:#fff;max-width:700px;margin:0 auto 20px auto;padding:40px 30px 40px 30px;border:1px solid #ccc;box-shadow:0 2px 8px #eee;font-family:Georgia,serif;line-height:1.7;">
                {{ html_output|safe }}
            </div>
            <h2>HTML à copier pour WordPress :</h2>
            <textarea readonly style="width:100%;height:200px;">{{ html_output }}</textarea>
            <h2>Insérer dans WordPress :</h2>
            <form method="post" action="/publish_wordpress">
                <textarea name="html_output" style="display:none;">{{ html_output }}</textarea>
                <input type="hidden" name="ideas" value="{{ ideas or '' }}">
                <input type="hidden" name="prompt" value="{{ prompt or '' }}">
                <input type="hidden" name="redesign_prompt" value="{{ redesign_prompt or '' }}">
                <input type="text" name="wp_title" placeholder="Titre de l'article" value="{{ ideas or '' }}" required style="width:100%;padding:8px;margin-bottom:8px;border-radius:6px;border:1px solid #bbb;">
                <input type="text" name="wp_categories" placeholder="Catégories WordPress (IDs, ex: 2,5)" style="width:100%;padding:8px;margin-bottom:8px;border-radius:6px;border:1px solid #bbb;">
                <input type="text" name="wp_tags" placeholder="Tags WordPress (IDs, ex: 12,18)" style="width:100%;padding:8px;margin-bottom:8px;border-radius:6px;border:1px solid #bbb;">
                <select name="wp_status" style="padding:8px;margin-bottom:8px;border-radius:6px;border:1px solid #bbb;">
                    <option value="draft">Brouillon</option>
                    <option value="publish">Publier directement</option>
                </select>
                <button type="submit" class="btn-main" style="margin-left:10px;">Insérer</button>
            </form>
            <form method="post" action="/test_wordpress" style="margin-top:8px;">
                <textarea name="html_output" style="display:none;">{{ html_output }}</textarea>
                <input type="hidden" name="ideas" value="{{ ideas or '' }}">
                <input type="hidden" name="prompt" value="{{ prompt or '' }}">
                <input type="hidden" name="redesign_prompt" value="{{ redesign_prompt or '' }}">
                <button type="submit" class="btn-main">Tester la connexion WordPress</button>
            </form>
            {% if wp_message %}<div style="margin-top:10px;padding:10px;border-radius:6px;background:{% if wp_success %}#d4edda{% else %}#f8d7da{% endif %};color:{% if wp_success %}#155724{% else %}#721c24{% endif %};">{{ wp_message }}</div>{% endif %}
            <h2>Instructions de modification du design/layout :</h2>
            <form method="post">
                <input type="hidden" name="ideas" value="{{ ideas }}">
                <input type="hidden" name="prompt" value="{{ prompt }}">
                <input type="hidden" name="html_output" value="{{ html_output }}">
                <textarea name="redesign_prompt" id="redesign_prompt" style="width:100%;height:80px;" placeholder="Ex: Mets les titres en bleu, ajoute une bordure..." required>{{ redesign_prompt or '' }}</textarea><br>
                <button type="submit" name="action" value="redesign" class="btn-redesign">Améliorer le design/layout</button>
            </form>
        </div>
        <script>
        // Persistance du champ redesign_prompt
        const redesignField = document.getElementById('redesign_prompt');
        if(localStorage.getItem('redesign_prompt')) redesignField.value = localStorage.getItem('redesign_prompt');
        redesignField.addEventListener('input', () => localStorage.setItem('redesign_prompt', redesignField.value));
        </script>
        {% endif %}
    </body>
    </html>

    '''


def is_logged_in():
    return session.get('logged_in', False)


from functools import wraps

PASSWORD = os.environ.get('APP_PASSWORD', 'Freeredac')  # Mot de passe personnalisé

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('logged_in'):
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated


# Route de connexion par mot de passe uniquement
@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    feedback_links = ''
    if request.method == 'POST':
        if request.form.get('password') == PASSWORD:
            session['logged_in'] = True
            return redirect(url_for('index'))
        else:
            error = 'Mot de passe incorrect.'
    return render_template_string(LOGIN_FORM, error=error)



@app.route('/', methods=['GET', 'POST'])
@login_required
def index():
    if not is_logged_in():
        return redirect(url_for('login'))
    html_output = None
    ideas = ''
    prompt = ''
    redesign_prompt = ''
    documents_content = ''
    links_content = ''
    links_extracted_text = ''
    feedback_links = ''
    generation_error = ''
    if request.method == 'POST':
        try:
            # Récupérer les fichiers joints
            if 'documents' in request.files:
                files = request.files.getlist('documents')
                for file in files:
                    if file.filename:

                        ext = os.path.splitext(file.filename)[1].lower()
                        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                            file.save(tmp.name)
                            tmp.seek(0)
                            if ext == '.txt':
                                documents_content += tmp.read().decode('utf-8', errors='ignore') + '\n'
                            elif ext == '.pdf':
                                try:
                                    tmp.close()  # Fermer avant lecture externe
                                    reader = PyPDF2.PdfReader(tmp.name)
                                    for page in reader.pages:
                                        documents_content += page.extract_text() + '\n'
                                except Exception as e:
                                    documents_content += f"[Erreur PDF: {e}]\n"
                            elif ext == '.docx':
                                try:
                                    tmp.close()
                                    doc = docx.Document(tmp.name)
                                    for para in doc.paragraphs:
                                        documents_content += para.text + '\n'
                                except Exception as e:
                                    documents_content += f"[Erreur DOCX: {e}]\n"
                            elif ext == '.pptx':
                                try:
                                    tmp.close()
                                    pres = pptx.Presentation(tmp.name)
                                    for slide in pres.slides:
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"):
                                                documents_content += shape.text + '\n'
                                except Exception as e:
                                    documents_content += f"[Erreur PPTX: {e}]\n"
                            elif ext in ['.xlsx', '.csv']:
                                try:
                                    tmp.close()
                                    if ext == '.csv':
                                        df = pd.read_csv(tmp.name)
                                    else:
                                        df = pd.read_excel(tmp.name)
                                    documents_content += df.to_string(index=False) + '\n'
                                except Exception as e:
                                    documents_content += f"[Erreur tableur: {e}]\n"
                            else:
                                tmp.close()
                            os.unlink(tmp.name)

            # Récupérer les liens web et extraire leur contenu
            links = request.form.get('links', '').strip()
            if links:
                links_content = links
                for link in links.splitlines():
                    url = link.strip()
                    if url:
                        try:
                            resp = requests.get(url, timeout=8)
                            resp.raise_for_status()
                            soup = BeautifulSoup(resp.text, 'html.parser')
                            for script in soup(['script', 'style', 'noscript']):
                                script.decompose()
                            text = soup.get_text(separator=' ', strip=True)
                            text = ' '.join(text.split())
                            links_extracted_text += f"\n[Contenu extrait de {url}]:\n{text[:3000]}\n"
                            feedback_links += f"<b>Succès extraction :</b> {url}<br>"
                        except Exception as e:
                            links_extracted_text += f"\n[Erreur lors de l'extraction de {url}: {e}]\n"
                            feedback_links += f"<b>Échec extraction :</b> {url} <span style='color:red;'>({e})</span><br>"

            # Si on demande un redesign
            if request.form.get('action') == 'redesign':
                redesign_prompt = request.form['redesign_prompt']
                html_output_old = request.form['html_output']
                ideas = request.form['ideas']
                prompt = request.form['prompt']
                redesign_full_prompt = (
                    f"Voici un article HTML :\n{html_output_old}\n\n{redesign_prompt} "
                    "Rends le HTML prêt pour WordPress, professionnel, académique, accessible, structuré, sans balises <html> ou <body>."
                )
                html_output = generer_article(ideas, redesign_full_prompt)
            else:
                ideas = request.form['ideas']
                prompt = request.form['prompt']
                redesign_prompt = ''
                # Combine tout pour le prompt IA
                prompt_html = prompt
                if documents_content:
                    prompt_html += f"\n\nVoici le contenu de documents joints :\n{documents_content}"
                if links_content:
                    prompt_html += f"\n\nVoici des liens à prendre en compte :\n{links_content}"
                if links_extracted_text:
                    prompt_html += f"\n\nVoici le contenu extrait des liens web fournis :\n{links_extracted_text}"
                prompt_html += (
                    "\n\nGénère un article au format HTML prêt à être inséré dans WordPress, avec : "
                    "<h2> pour les titres principaux, <h3> pour les sous-titres, <p> pour les paragraphes, <ul> pour les listes, <blockquote> pour les citations. "
                    "Commence par une introduction claire, structure l'article avec des titres explicites, des paragraphes courts, des exemples concrets, des citations si pertinent, et termine par une conclusion synthétique. "
                    "Le style doit être professionnel, académique, intelligent mais accessible à tous. N'utilise pas de balises <html> ou <body>."
                )
                html_output = generer_article(ideas, prompt_html)
        except Exception as e:
            generation_error = f"Erreur pendant la génération: {e}"
    return render_template_string(
        HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
        html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links=feedback_links,
        wp_message=None, wp_success=False, generation_error=generation_error
    )

@app.route('/publish_wordpress', methods=['POST'])
@login_required
def publish_wordpress():
    html_output = request.form.get('html_output', '')
    title = request.form.get('wp_title', 'Article généré par IA').strip()
    status = request.form.get('wp_status', 'draft')
    wp_categories = parse_id_list(request.form.get('wp_categories', ''))
    wp_tags = parse_id_list(request.form.get('wp_tags', ''))
    ideas = request.form.get('ideas', '')
    prompt = request.form.get('prompt', '')
    redesign_prompt = request.form.get('redesign_prompt', '')

    if not html_output.strip():
        return render_template_string(
            HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
            html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
            wp_message='Aucun contenu HTML à insérer dans WordPress.', wp_success=False
        )

    if not title:
        title = 'Article généré par IA'

    wp_url, wp_user, wp_password = get_wp_auth()
    post_type = get_wp_post_type()

    if not wp_url or not wp_user or not wp_password:
        return render_template_string(
            HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
            html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
            wp_message='Variables WP_URL, WP_USERNAME ou WP_APP_PASSWORD manquantes dans Render Environment.', wp_success=False
        )

    headers = {'Content-Type': 'application/json'}
    payload = {'title': title, 'content': html_output, 'status': status}
    if wp_categories:
        payload['categories'] = wp_categories
    if wp_tags:
        payload['tags'] = wp_tags

    try:
        resp = wp_request('POST', f"{wp_url}/wp-json/wp/v2/{post_type}", wp_user, wp_password, json=payload, headers=headers, timeout=15)
        if resp.status_code in (200, 201):
            post_url = resp.json().get('link', '')
            msg = f'Article inséré avec succès ! <a href="{post_url}" target="_blank">Voir l\'article</a>'
            return render_template_string(
                HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
                html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
                wp_message=msg, wp_success=True
            )
        else:
            return render_template_string(
                HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
                html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
                wp_message=build_wp_error_message(resp), wp_success=False
            )
    except Exception as e:
        return render_template_string(
            HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
            html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
            wp_message=f'Erreur de connexion WordPress: {e}', wp_success=False
        )


@app.route('/test_wordpress', methods=['POST'])
@login_required
def test_wordpress():
    html_output = request.form.get('html_output', '')
    ideas = request.form.get('ideas', '')
    prompt = request.form.get('prompt', '')
    redesign_prompt = request.form.get('redesign_prompt', '')

    wp_url, wp_user, wp_password = get_wp_auth()
    post_type = get_wp_post_type()
    if not wp_url or not wp_user or not wp_password:
        return render_template_string(
            HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
            html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
            wp_message='Variables WP_URL, WP_USERNAME ou WP_APP_PASSWORD manquantes dans Render Environment.', wp_success=False
        )

    try:
        preflight = requests.get(f"{wp_url}/wp-json/", timeout=12)
        if preflight.status_code >= 500:
            message = (
                "Le endpoint public /wp-json renvoie déjà une erreur serveur. "
                + build_wp_error_message(preflight)
                + " Vérifiez les logs PHP/Apache et testez temporairement sans plugins de sécurité/cache."
            )
            success = False
            return render_template_string(
                HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
                html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
                wp_message=message, wp_success=success
            )

        resp = wp_request('GET', f"{wp_url}/wp-json/wp/v2/users/me?context=edit", wp_user, wp_password, timeout=12)
        if resp.status_code == 200:
            user_data = resp.json()
            user_name = user_data.get('name') or wp_user
            roles = ', '.join(user_data.get('roles', [])) or 'inconnu'
            message = (
                f'Connexion WordPress OK. Utilisateur authentifié: {user_name}. '
                f'Rôles détectés: {roles}. '
                f"Type de contenu ciblé: {post_type}. "
                "Si l'insertion échoue avec rest_cannot_create, ce compte n'a pas la capacité de créer ce type de contenu."
            )
            success = True
        else:
            message = build_wp_error_message(resp)
            success = False
    except Exception as e:
        message = f'Erreur de connexion WordPress: {e}'
        success = False

    return render_template_string(
        HTML_FORM + '<br><br><a href="/logout">Se déconnecter</a>',
        html_output=html_output, ideas=ideas, prompt=prompt, redesign_prompt=redesign_prompt, feedback_links='',
        wp_message=message, wp_success=success
    )


if __name__ == '__main__':
    print('Flask démarre...')
    app.run(host='0.0.0.0', port=5000, debug=True)
